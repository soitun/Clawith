"""Editable PPTX rendering implementation for HTML inputs."""

import re
from pathlib import Path
from typing import Any

from loguru import logger

from app.services.document_conversion.chrome_renderer import collect_browser_layout


async def render_html_to_pptx(src_file: Path, tgt_file: Path, target_path: str, ws: Path, arguments: dict[str, Any]) -> str:
    try:
        from bs4 import BeautifulSoup
        from bs4.element import Tag
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
        
        html_content = src_file.read_text(encoding="utf-8")
        soup = BeautifulSoup(html_content, "html.parser")

        design_w_px = int(arguments.get("design_width") or 1280)
        design_h_px = int(arguments.get("design_height") or 720)
        render_mode = str(arguments.get("render_mode") or "editable").lower()
        try:
            render_scale = float(arguments.get("render_scale") or 2.0)
        except (TypeError, ValueError):
            render_scale = 2.0
        render_scale = max(1.0, min(4.0, render_scale))
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        named_colors = {
            "black": "000000", "white": "ffffff", "gray": "808080", "grey": "808080",
            "red": "ff0000", "green": "008000", "blue": "0000ff", "transparent": "",
        }

        def parse_css_block(css: str) -> dict[str, dict[str, str]]:
            rules: dict[str, dict[str, str]] = {}
            css = re.sub(r"/\*.*?\*/", "", css, flags=re.S)
            for selector_text, body in re.findall(r"([^{}]+)\{([^{}]+)\}", css):
                decls = parse_style(body)
                for selector in selector_text.split(","):
                    selector = selector.strip()
                    if selector:
                        rules.setdefault(selector, {}).update(decls)
            return rules

        def parse_style(style: str | None) -> dict[str, str]:
            out: dict[str, str] = {}
            if not style:
                return out
            for part in style.split(";"):
                if ":" not in part:
                    continue
                key, value = part.split(":", 1)
                key = key.strip().lower()
                value = value.strip()
                if key and value:
                    out[key] = value
            return out

        css_rules: dict[str, dict[str, str]] = {}
        for style_tag in soup.find_all("style"):
            css_rules.update(parse_css_block(style_tag.get_text("\n")))

        def element_style(el: Tag | None) -> dict[str, str]:
            if not isinstance(el, Tag):
                return {}
            style: dict[str, str] = {}
            for selector in ("html", "body"):
                style.update(css_rules.get(selector, {}))
            style.update(css_rules.get(el.name or "", {}))
            for cls in el.get("class") or []:
                style.update(css_rules.get(f".{cls}", {}))
                style.update(css_rules.get(f"{el.name}.{cls}", {}))
            if el.get("id"):
                style.update(css_rules.get(f"#{el.get('id')}", {}))
            style.update(parse_style(el.get("style")))
            return style

        def color_tuple(value: str | None) -> tuple[int, int, int, float] | None:
            if not value:
                return None
            value = value.strip().lower()
            color_match = re.search(r"#(?:[0-9a-f]{3}|[0-9a-f]{6})\b|rgba?\([^)]+\)", value)
            if color_match:
                value = color_match.group(0)
            if value in named_colors:
                value = named_colors[value]
            if not value or value in ("none", "transparent"):
                return None
            alpha_match = re.match(r"rgba\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([0-9.]+)\s*\)", value)
            if alpha_match:
                try:
                    if float(alpha_match.group(1)) <= 0.01:
                        return None
                except ValueError:
                    return None
            if value.startswith("#"):
                raw = value[1:]
                if len(raw) == 3:
                    raw = "".join(ch * 2 for ch in raw)
                if len(raw) == 6:
                    try:
                        return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16), 1.0)
                    except ValueError:
                        return None
            m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([0-9.]+))?", value)
            if m:
                alpha = 1.0
                if m.group(4) is not None:
                    try:
                        alpha = max(0.0, min(1.0, float(m.group(4))))
                    except ValueError:
                        alpha = 1.0
                if alpha <= 0.01:
                    return None
                return (int(m.group(1)), int(m.group(2)), int(m.group(3)), alpha)
            return None

        def parse_color(value: str | None, backdrop: str | RGBColor | None = None) -> RGBColor | None:
            rgba = color_tuple(value)
            if not rgba:
                return None
            r, g, b, alpha = rgba
            if alpha < 0.999:
                backdrop_tuple: tuple[int, int, int, float] | None = None
                if isinstance(backdrop, RGBColor):
                    backdrop_tuple = (int(backdrop[0]), int(backdrop[1]), int(backdrop[2]), 1.0)
                elif isinstance(backdrop, str):
                    backdrop_tuple = color_tuple(backdrop)
                if backdrop_tuple:
                    br, bg, bb, _ = backdrop_tuple
                    r = round(r * alpha + br * (1 - alpha))
                    g = round(g * alpha + bg * (1 - alpha))
                    b = round(b * alpha + bb * (1 - alpha))
            return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))

        def representative_color(value: str | None, prefer: str = "last", backdrop: str | RGBColor | None = None) -> RGBColor | None:
            if not value:
                return None
            matches = re.findall(r"#(?:[0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b|rgba?\([^)]+\)", value)
            if not matches:
                return parse_color(value, backdrop)
            chosen = matches[-1] if prefer == "last" else matches[0]
            return parse_color(chosen, backdrop)

        def length_to_inches(value: str | None, axis_total_in: float, axis_px: int) -> float | None:
            if not value:
                return None
            value = str(value).strip().lower()
            try:
                if value.endswith("%"):
                    return axis_total_in * float(value[:-1]) / 100.0
                if value.endswith("px"):
                    return axis_total_in * float(value[:-2]) / axis_px
                if value.endswith("pt"):
                    return float(value[:-2]) / 72.0
                if value.endswith("in"):
                    return float(value[:-2])
                if value.endswith("rem") or value.endswith("em"):
                    return axis_total_in * (float(value[:-3] if value.endswith("rem") else value[:-2]) * 16) / axis_px
                return axis_total_in * float(value) / axis_px
            except ValueError:
                return None

        def font_size_pt(style: dict[str, str], default: int) -> float:
            raw = style.get("font-size")
            if not raw:
                return float(default)
            raw = raw.strip().lower()
            try:
                if raw.endswith("px"):
                    return float(raw[:-2]) * 0.75
                if raw.endswith("pt"):
                    return float(raw[:-2])
                if raw.endswith("rem") or raw.endswith("em"):
                    return float(raw[:-3] if raw.endswith("rem") else raw[:-2]) * 12
            except ValueError:
                return float(default)
            return float(default)

        def text_align(style: dict[str, str]):
            align = (style.get("text-align") or "").lower()
            return {
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }.get(align, PP_ALIGN.LEFT)

        def vertical_anchor(style: dict[str, str]):
            display = (style.get("display") or "").lower()
            align_items = (style.get("align-items") or "").lower()
            justify_content = (style.get("justify-content") or "").lower()
            if "flex" in display and ("center" in align_items or "center" in justify_content):
                return MSO_ANCHOR.MIDDLE
            return MSO_ANCHOR.TOP

        def css_px_to_inches(value: str | None, axis_px: int = 1280) -> float:
            raw = str(value or "").strip().lower()
            if not raw or raw in ("auto", "normal"):
                return 0.0
            try:
                if raw.endswith("px"):
                    return float(raw[:-2]) * (prs.slide_width / 914400) / axis_px
                if raw.endswith("rem"):
                    return float(raw[:-3]) * 16 * (prs.slide_width / 914400) / axis_px
                if raw.endswith("em"):
                    return float(raw[:-2]) * 16 * (prs.slide_width / 914400) / axis_px
            except ValueError:
                return 0.0
            return 0.0

        def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, style: dict[str, str], default_size: int, bold: bool = False):
            shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = vertical_anchor(style)
            tf.margin_left = Inches(min(css_px_to_inches(style.get("padding-left"), design_w_px), max(w * 0.35, 0)))
            tf.margin_right = Inches(min(css_px_to_inches(style.get("padding-right"), design_w_px), max(w * 0.35, 0)))
            tf.margin_top = Inches(min(css_px_to_inches(style.get("padding-top"), design_w_px), max(h * 0.35, 0)))
            tf.margin_bottom = Inches(min(css_px_to_inches(style.get("padding-bottom"), design_w_px), max(h * 0.35, 0)))
            p = tf.paragraphs[0]
            p.alignment = text_align(style)
            run = p.add_run()
            run.text = text
            font = run.font
            font.size = Pt(font_size_pt(style, default_size))
            font.bold = bold or (style.get("font-weight") or "").lower() in ("bold", "600", "700", "800", "900")
            color = parse_color(style.get("color"), style.get("_backdrop-color"))
            text_fill = (style.get("-webkit-text-fill-color") or "").lower()
            bg_clip = (style.get("background-clip") or style.get("-webkit-background-clip") or "").lower()
            if text_fill in ("transparent", "rgba(0, 0, 0, 0)") or "text" in bg_clip:
                color = representative_color(style.get("background"), "last", style.get("_backdrop-color")) or color
            if color:
                font.color.rgb = color
            family = (style.get("font-family") or "").split(",")[0].strip().strip("\"'")
            if family:
                font.name = family
            return shape

        def add_background(slide, style: dict[str, str]):
            color = parse_color(style.get("background") or style.get("background-color"))
            if color:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = color

        def add_card(slide, x: float, y: float, w: float, h: float, style: dict[str, str]):
            backdrop = style.get("_backdrop-color")
            bg = parse_color(style.get("background") or style.get("background-color"), backdrop)
            border = parse_color(style.get("border-color") or style.get("border"), bg or backdrop)
            has_border = "border" in style
            try:
                has_border = has_border or float(str(style.get("border-width") or "0").replace("px", "").strip() or 0) > 0
            except ValueError:
                pass
            if not bg and has_border:
                border = border or RGBColor(220, 224, 232)
            if not has_border:
                border = None
            if not bg and not border:
                return None
            radius = length_to_inches(style.get("border-radius"), prs.slide_width / 914400, design_w_px) or 0
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0.03 else MSO_SHAPE.RECTANGLE
            shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
            if bg:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg
            else:
                shape.fill.background()
            if border:
                shape.line.color.rgb = border
            else:
                shape.line.color.rgb = bg or RGBColor(220, 224, 232)
            return shape

        def image_path(src: str | None) -> Path | None:
            if not src or src.startswith(("http://", "https://", "data:")):
                return None
            if src.startswith("file://"):
                from urllib.parse import unquote, urlparse

                p = Path(unquote(urlparse(src).path))
                return p if p.exists() else None
            p = Path(src)
            candidates = [src_file.parent / p, ws / p]
            for candidate in candidates:
                if candidate.exists():
                    return candidate
            return None

        def element_box(style: dict[str, str]) -> tuple[float | None, float | None, float | None, float | None]:
            sw = prs.slide_width / 914400
            sh = prs.slide_height / 914400
            return (
                length_to_inches(style.get("left") or style.get("x"), sw, design_w_px),
                length_to_inches(style.get("top") or style.get("y"), sh, design_h_px),
                length_to_inches(style.get("width") or style.get("max-width"), sw, design_w_px),
                length_to_inches(style.get("height") or style.get("min-height"), sh, design_h_px),
            )

        def visible_children(el: Tag) -> list[Tag]:
            return [child for child in el.children if isinstance(child, Tag) and child.name not in ("style", "script", "meta", "link")]

        def render_flow_element(slide, el: Tag, y: float, x: float = 0.75, width: float = 11.85) -> float:
            style = element_style(el)
            name = el.name or ""
            left, top, box_w, box_h = element_box(style)
            if (style.get("position") == "absolute" or left is not None or top is not None) and (left is not None or top is not None):
                render_absolute_element(slide, el, left or x, top or y, box_w or width, box_h)
                return y
            if name == "img":
                p = image_path(el.get("src"))
                if p:
                    h = box_h or 2.2
                    slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(box_w or min(width, 5.5)), height=Inches(h))
                    return y + h + 0.18
                return y
            classes = set(el.get("class") or [])
            looks_like_card = bool({"card", "panel", "box", "tile"} & classes) or any(k in style for k in ("background", "background-color", "border", "border-color"))
            children = visible_children(el)
            if children and name in ("div", "main", "section", "article", "header", "footer", "aside", "nav"):
                if looks_like_card:
                    text_len = len(el.get_text(" ", strip=True))
                    h = box_h or max(1.0, min(4.8, 0.42 * max(text_len // 42 + 1, 2)))
                    add_card(slide, x, y, box_w or width, h, style)
                    inner_y = y + 0.18
                    for child in children:
                        inner_y = render_flow_element(slide, child, inner_y, x + 0.25, (box_w or width) - 0.5)
                    return y + max(h, inner_y - y + 0.18) + 0.18
                for child in children:
                    y = render_flow_element(slide, child, y, x, width)
                return y
            text = el.get_text(" ", strip=True)
            if not text:
                for child in visible_children(el):
                    y = render_flow_element(slide, child, y, x, width)
                return y
            if name in ("h1",):
                add_textbox(slide, text, x, y, width, box_h or 0.85, style, 40, True)
                return y + (box_h or 0.95)
            if name in ("h2",):
                add_textbox(slide, text, x, y, width, box_h or 0.62, style, 30, True)
                return y + (box_h or 0.74)
            if name in ("h3",):
                add_textbox(slide, text, x, y, width, box_h or 0.45, style, 22, True)
                return y + (box_h or 0.55)
            if name in ("li",):
                add_textbox(slide, f"• {text}", x + 0.2, y, width - 0.2, box_h or 0.36, style, 18)
                return y + (box_h or 0.42)
            if name in ("ul", "ol"):
                for li in el.find_all("li", recursive=False):
                    y = render_flow_element(slide, li, y, x, width)
                return y

            if looks_like_card:
                h = box_h or max(1.0, min(3.2, 0.42 * max(len(text) // 42 + 1, 2)))
                add_card(slide, x, y, box_w or width, h, style)
                inner_y = y + 0.18
                if children:
                    for child in children:
                        inner_y = render_flow_element(slide, child, inner_y, x + 0.25, (box_w or width) - 0.5)
                else:
                    add_textbox(slide, text, x + 0.25, inner_y, (box_w or width) - 0.5, h - 0.3, style, 18)
                return y + h + 0.18

            add_textbox(slide, text, x, y, width, box_h or 0.42, style, 18)
            return y + (box_h or 0.5)

        def render_absolute_element(slide, el: Tag, x: float, y: float, w: float, h: float | None):
            style = element_style(el)
            name = el.name or ""
            if name == "img":
                p = image_path(el.get("src"))
                if p:
                    slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h or 2.0))
                return
            if name in ("div", "section", "article", "main"):
                add_card(slide, x, y, w, h or 1.5, style)
                children = visible_children(el)
                if children:
                    inner_y = y + 0.15
                    for child in children:
                        inner_y = render_flow_element(slide, child, inner_y, x + 0.2, max(w - 0.4, 0.5))
                    return
            text = el.get_text(" ", strip=True)
            if text:
                default = 38 if name == "h1" else 28 if name == "h2" else 18
                add_textbox(slide, text, x, y, w, h or 0.6, style, default, name in ("h1", "h2", "h3"))

        def render_browser_layout(layout: dict[str, Any]) -> bool:
            slides = layout.get("slides") or []
            if not slides:
                return False
            slide_w = prs.slide_width / 914400
            slide_h = prs.slide_height / 914400

            for slide_data in slides:
                slide_bg_value = slide_data.get("backgroundColor") or ""
                root_w = float(slide_data.get("width") or design_w_px or 1280)
                root_h = float(slide_data.get("height") or design_h_px or 720)
                preserve_aspect = bool(slide_data.get("backgroundScreenshots")) or root_h < root_w * 0.45
                if preserve_aspect:
                    scale = min(slide_w / root_w, slide_h / root_h)
                    sx = sy = scale
                    offset_x = (slide_w - root_w * scale) / 2
                    offset_y = (slide_h - root_h * scale) / 2
                else:
                    sx = slide_w / root_w
                    sy = slide_h / root_h
                    offset_x = 0.0
                    offset_y = 0.0
                slide = prs.slides.add_slide(blank_layout)
                add_background(slide, {"background-color": slide_bg_value})
                bg_screenshots = layout.get("backgroundScreenshots") or []
                bg_screenshot = bg_screenshots[len(prs.slides) - 1] if len(bg_screenshots) >= len(prs.slides) else None
                if bg_screenshot and Path(bg_screenshot).exists():
                    slide.shapes.add_picture(
                        bg_screenshot,
                        Inches(offset_x),
                        Inches(offset_y),
                        width=Inches(root_w * sx),
                        height=Inches(root_h * sy),
                    )

                for item in slide_data.get("items") or []:
                    style = item.get("style") or {}
                    raw_x = float(item.get("x") or 0)
                    raw_y = float(item.get("y") or 0)
                    raw_w = float(item.get("w") or 1)
                    raw_h = float(item.get("h") or 1)
                    if raw_x >= root_w or raw_y >= root_h or raw_x + raw_w <= 0 or raw_y + raw_h <= 0:
                        continue
                    x = max(0.0, offset_x + float(item.get("x") or 0) * sx)
                    y = max(0.0, offset_y + float(item.get("y") or 0) * sy)
                    w = max(0.05, min(raw_w, max(1.0, root_w - raw_x)) * sx)
                    h = max(0.05, min(raw_h, max(1.0, root_h - raw_y)) * sy)
                    ppt_style = {
                        "background": "" if style.get("backgroundImage") == "none" else (style.get("backgroundImage") or ""),
                        "background-color": style.get("backgroundColor") or "",
                        "border-color": style.get("borderColor") or "",
                        "border-width": style.get("borderWidth") or "",
                        "border-radius": style.get("borderRadius") or "",
                        "color": style.get("color") or "",
                        "font-size": style.get("fontSize") or "",
                        "font-family": style.get("fontFamily") or "",
                        "font-weight": style.get("fontWeight") or "",
                        "text-align": style.get("textAlign") or "",
                        "display": style.get("display") or "",
                        "align-items": style.get("alignItems") or "",
                        "justify-content": style.get("justifyContent") or "",
                        "padding-left": style.get("paddingLeft") or "",
                        "padding-right": style.get("paddingRight") or "",
                        "padding-top": style.get("paddingTop") or "",
                        "padding-bottom": style.get("paddingBottom") or "",
                        "-webkit-text-fill-color": style.get("webkitTextFillColor") or "",
                        "background-clip": style.get("backgroundClip") or "",
                        "-webkit-background-clip": style.get("webkitBackgroundClip") or "",
                        "_backdrop-color": slide_bg_value,
                    }
                    kind = item.get("kind")
                    if kind == "shape":
                        shape_screenshots = layout.get("shapeScreenshots") or {}
                        shape_screenshot = shape_screenshots.get(str(item.get("itemId") or ""))
                        if shape_screenshot and Path(shape_screenshot).exists():
                            slide.shapes.add_picture(shape_screenshot, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
                        else:
                            add_card(slide, x, y, w, h, ppt_style)
                    elif kind == "image":
                        src = item.get("src") or ""
                        p = image_path(src)
                        if p:
                            slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
                    elif kind == "text":
                        text = str(item.get("text") or "").strip()
                        if not text:
                            continue
                        tag = item.get("tag") or ""
                        default = 38 if tag == "h1" else 28 if tag == "h2" else 22 if tag == "h3" else 18
                        add_textbox(slide, text, x, y, w, h, ppt_style, default, tag in ("h1", "h2", "h3", "strong"))
            return True

        def render_browser_screenshots(layout: dict[str, Any]) -> bool:
            slides = layout.get("slides") or []
            screenshots = layout.get("screenshots") or []
            if not slides or not screenshots:
                return False
            slide_w = prs.slide_width / 914400
            slide_h = prs.slide_height / 914400

            for slide_data, screenshot in zip(slides, screenshots):
                if not screenshot or not Path(screenshot).exists():
                    return False
                slide = prs.slides.add_slide(blank_layout)
                add_background(slide, {"background-color": slide_data.get("backgroundColor") or ""})
                root_w = max(1.0, float(slide_data.get("width") or design_w_px))
                root_h = max(1.0, float(slide_data.get("height") or design_h_px))
                scale = min(slide_w / root_w, slide_h / root_h)
                image_w = root_w * scale
                image_h = root_h * scale
                left = (slide_w - image_w) / 2
                top = (slide_h - image_h) / 2
                slide.shapes.add_picture(
                    screenshot,
                    Inches(left),
                    Inches(top),
                    width=Inches(image_w),
                    height=Inches(image_h),
                )
            return True

        browser_layout = await collect_browser_layout(src_file, design_w_px, design_h_px, render_mode, render_scale)
        if browser_layout and render_mode in ("visual", "screenshot", "image", "hybrid") and render_browser_screenshots(browser_layout):
            tgt_file.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(tgt_file))
            return (
                f"✅ Successfully converted HTML to high-fidelity PPTX screenshots: {target_path}\n"
                "Note: visual style is preserved by rendering each slide from Chrome. Text and layout are not directly editable; "
                "use render_mode='editable' when editable PPT elements are more important than visual fidelity."
            )

        if browser_layout and render_browser_layout(browser_layout):
            tgt_file.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(tgt_file))
            return (
                f"✅ Successfully converted HTML to editable PPTX with browser layout: {target_path}\n"
                "Note: positions, sizes, colors, typography, cards, lists, and local images are sampled from a real browser. "
                "Effects such as shadows, filters, gradients, and complex text wrapping may still differ from HTML."
            )

        body = soup.body or soup
        slide_nodes = soup.select(".slide")
        if not slide_nodes:
            slide_nodes = [node for node in body.find_all(["section", "article"], recursive=False)]
        if not slide_nodes:
            slide_nodes = [body]

        for slide_node in slide_nodes:
            if not isinstance(slide_node, Tag):
                continue
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, element_style(slide_node) or element_style(body))
            current_y = 0.65
            children = visible_children(slide_node)
            if not children and slide_node is not body:
                text = slide_node.get_text(" ", strip=True)
                if text:
                    add_textbox(slide, text, 0.85, current_y, 11.6, 5.8, element_style(slide_node), 24)
                continue
            for child in children:
                current_y = render_flow_element(slide, child, current_y)
                if current_y > 7.0:
                    break
                    
        tgt_file.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(tgt_file))
        return (
            f"✅ Successfully converted HTML to editable PPTX: {target_path}\n"
            "Note: common typography, colors, cards, lists, images, and simple absolute positioning are preserved; "
            "complex CSS such as flex/grid effects, shadows, filters, and animations may still need manual adjustment."
        )
    except Exception as e:
        logger.exception(f"Convert HTML to PPTX failed: {e}")
        return f"❌ Conversion failed: {e}"
